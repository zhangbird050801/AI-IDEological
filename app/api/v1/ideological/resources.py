from typing import List
from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File, Form
from fastapi.responses import StreamingResponse
from tortoise.queryset import Q
from tortoise.expressions import F
import os
import uuid
import aiofiles
import httpx
from pathlib import Path
import json
import re
import tempfile
import zipfile
from starlette.concurrency import run_in_threadpool
from app.schemas.ideological import (
    TeachingResourceCreate,
    TeachingResourceUpdate,
    TeachingResource,
    ResourceSearchRequest,
    BatchOperationRequest,
    BatchOperationResponse,
)
from app.models.ideological import TeachingResource as TeachingResourceModel, ResourceType
from app.models.chapter import Chapter
from app.models.admin import User
from app.core.dependency import AuthControl
from app.core.crud import CRUDBase
from app.core.aigc.aigc_client import AIGCClient
from app.settings.config import settings
from app.services.recommendation_service import RecommendationService

router = APIRouter()

# 配置文件上传目录（相对项目根）
UPLOAD_DIR = Path("uploads/teaching_resources")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)

class ResourceService(CRUDBase[TeachingResourceModel, TeachingResourceCreate, TeachingResourceUpdate]):
    def __init__(self):
        super().__init__(TeachingResourceModel)

    async def create_resource(self, obj_in: TeachingResourceCreate, user_id: int, file_path: str = None) -> TeachingResourceModel:
        obj_data = obj_in.dict()
        obj_data["uploader_id"] = user_id
        if file_path:
            obj_data["file_path"] = file_path
        obj_data = await _hydrate_course_chapter(obj_data)
        return await self.create(obj_data)

    async def get_resources_with_search(self, search_request: ResourceSearchRequest, user_id: int = None):
        query = TeachingResourceModel.all()

        # 如果不是获取所有公开资源，需要根据用户权限过滤
        if user_id:
            query = query.filter(
                Q(is_public=True) | Q(uploader_id=user_id)
            )
        else:
            query = query.filter(is_public=True)

        # 关键词搜索
        if search_request.keyword:
            keyword = search_request.keyword.strip()
            query = query.filter(
                Q(title__icontains=keyword) |
                Q(description__icontains=keyword)
            )

        # 其他筛选条件
        if search_request.resource_type:
            query = query.filter(resource_type=search_request.resource_type)

        # 课程和章节筛选
        # 逻辑：course_id 和 chapter_id 应该是 AND 关系
        if search_request.course_id:
            query = query.filter(course_id=search_request.course_id)
        
        if search_request.chapter_id:
            query = query.filter(chapter_id=search_request.chapter_id)
        
        if search_request.software_engineering_chapter:
            query = query.filter(software_engineering_chapter__icontains=search_request.software_engineering_chapter)

        if search_request.theme_category_id:
            query = query.filter(theme_category_id=search_request.theme_category_id)

        if search_request.tags:
            for tag in search_request.tags:
                query = query.filter(tags__contains=tag)

        # 排序：按创建时间倒序，按使用次数倒序
        query = query.order_by("-created_at", "-usage_count")

        # 分页
        offset = (search_request.page - 1) * search_request.page_size
        total = await query.count()
        items = await query.offset(offset).limit(search_request.page_size).prefetch_related('uploader')

        return {
            "items": items,
            "total": total,
            "page": search_request.page,
            "page_size": search_request.page_size,
            "pages": (total + search_request.page_size - 1) // search_request.page_size
        }

    async def update_resource_usage(self, resource_id: int):
        resource = await TeachingResourceModel.get_or_none(id=resource_id)
        if not resource:
            raise HTTPException(status_code=404, detail="资源不存在")

        await resource.update_from_dict({"usage_count": resource.usage_count + 1})
        await resource.save()

        return resource

    async def get_hot_resources(self, limit: int = 10):
        return await TeachingResourceModel.filter(
            is_public=True
        ).order_by("-usage_count", "-created_at").limit(limit)

    async def upload_file(self, file: UploadFile, user_id: int) -> dict:
        # 验证文件类型
        allowed_extensions = {
            "document": [".pdf", ".doc", ".docx", ".txt", ".rtf"],
            "video": [".mp4", ".avi", ".mov", ".wmv", ".flv"],
            "audio": [".mp3", ".wav", ".flac", ".aac"],
            "image": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".svg"],
            "presentation": [".ppt", ".pptx", ".key"],
            "other": []
        }

        file_extension = Path(file.filename).suffix.lower()

        # 确定文件类型
        resource_type = "other"
        for type_name, extensions in allowed_extensions.items():
            if file_extension in extensions:
                resource_type = type_name
                break

        # 生成唯一文件名
        file_uuid = uuid.uuid4()
        file_name = f"{file_uuid}{file_extension}"
        file_path = UPLOAD_DIR / file_name

        # 保存文件
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)

        # 获取文件大小
        file_size = len(content)

        # 构建访问URL
        # 静态可访问的URL（需在 FastAPI 挂载 /uploads）
        static_file_url = f"/uploads/teaching_resources/{file_name}"
        download_url = f"/api/v1/ideological/resources/download/{file_uuid}"

        # 预览URL：图片/PDF/DOCX 直接用静态地址，方便前端内嵌预览
        preview_url = static_file_url if resource_type in ["image", "document"] else None

        return {
            "file_path": str(file_path),
            "file_url": static_file_url,
            "download_url": download_url,
            "preview_url": preview_url,
            "file_size": file_size,
            "file_format": file_extension.lstrip('.'),
            "resource_type": resource_type
        }


async def _hydrate_course_chapter(obj_data: dict) -> dict:
    chapter_id = obj_data.get("chapter_id")
    course_id = obj_data.get("course_id")
    chapter_name = obj_data.get("software_engineering_chapter")

    if chapter_id:
        chapter = await Chapter.get_or_none(id=chapter_id)
        if chapter:
            obj_data["software_engineering_chapter"] = chapter.name
            if not course_id:
                obj_data["course_id"] = chapter.course_id
        return obj_data

    if chapter_name:
        query = Chapter.filter(name=chapter_name)
        if course_id:
            query = query.filter(course_id=course_id)
        chapter = await query.first()
        if chapter:
            obj_data["chapter_id"] = chapter.id
            if not course_id:
                obj_data["course_id"] = chapter.course_id

    return obj_data

def _normalize_text(text: str) -> str:
    if not text:
        return ''
    return '\n'.join(line.rstrip() for line in text.splitlines()).strip()

def _normalize_kimi_content(raw_text: str) -> str:
    text = raw_text or ''
    image_pattern = r'!\[[^\]]*]\([^)]+\)(\{[^}]+\})?'
    try:
        data = json.loads(text)
        if isinstance(data, dict):
            content = (
                data.get("content")
                or data.get("text")
                or (data.get("data", {}) or {}).get("content")
                or (data.get("data", {}) or {}).get("text")
            )
            if isinstance(content, str):
                text = content
    except Exception:
        pass

    if re.search(image_pattern, text):
        text = re.sub(image_pattern, '', text)

    text = _normalize_text(text)
    if not text and raw_text and re.search(image_pattern, raw_text):
        return "（该文件主要为图片或扫描件，未抽取到可用文本）"
    return text

def _needs_image_fallback(text: str) -> bool:
    if not text:
        return True
    if text.strip() == "（该文件主要为图片或扫描件，未抽取到可用文本）":
        return True
    return False

async def _extract_docx_images_text(file_path: Path, client: AIGCClient) -> str:
    if not file_path.exists():
        return ''
    if file_path.suffix.lower() != '.docx':
        return ''

    image_texts = []
    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            with zipfile.ZipFile(str(file_path)) as zf:
                image_names = [n for n in zf.namelist() if n.startswith("word/media/")]
                for name in image_names:
                    data = zf.read(name)
                    if not data:
                        continue
                    image_path = Path(tmp_dir) / Path(name).name
                    image_path.write_bytes(data)
                    raw_text = await client.extract_file_text(str(image_path))
                    normalized = _normalize_kimi_content(raw_text)
                    if normalized and not _needs_image_fallback(normalized):
                        image_texts.append(normalized)
        except Exception:
            return ''

    return _normalize_text('\n'.join(image_texts))

def _extract_text_sync(file_path: Path, file_format: str, resource_type: str) -> str:
    suffix = (file_format or '').lower().lstrip('.')
    if not suffix:
        suffix = file_path.suffix.lower().lstrip('.')

    if suffix in ['txt', 'rtf']:
        for encoding in ['utf-8', 'utf-16', 'gbk', 'latin-1']:
            try:
                return file_path.read_text(encoding=encoding)
            except Exception:
                continue
        return file_path.read_text(errors='ignore')

    if suffix == 'pdf':
        try:
            import pdfplumber
        except Exception as exc:
            raise ValueError("缺少依赖 pdfplumber，请安装后重试") from exc
        text_parts = []
        with pdfplumber.open(str(file_path)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ''
                if page_text:
                    text_parts.append(page_text)
        return '\n'.join(text_parts)

    if suffix in ['docx', 'doc']:
        if suffix == 'doc':
            raise ValueError('暂不支持 .doc 格式，请转换为 .docx')
        try:
            import docx
        except Exception as exc:
            raise ValueError("缺少依赖 python-docx，请安装后重试") from exc
        doc = docx.Document(str(file_path))
        return '\n'.join(p.text for p in doc.paragraphs if p.text)

    if suffix in ['pptx', 'ppt']:
        if suffix == 'ppt':
            raise ValueError('暂不支持 .ppt 格式，请转换为 .pptx')
        try:
            from pptx import Presentation
        except Exception as exc:
            raise ValueError("缺少依赖 python-pptx，请安装: pip install python-pptx") from exc
        
        text_parts = []
        prs = Presentation(str(file_path))
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            # 提取幻灯片中的所有文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
            if slide_texts:
                text_parts.append(f"--- 幻灯片 {slide_num} ---\n" + '\n'.join(slide_texts))
        return '\n\n'.join(text_parts)

    if resource_type == 'image' or suffix in ['png', 'jpg', 'jpeg', 'bmp', 'gif', 'tiff']:
        try:
            from PIL import Image
        except Exception as exc:
            raise ValueError("缺少依赖 pillow，请安装后重试") from exc
        try:
            import pytesseract
        except Exception as exc:
            raise ValueError("缺少依赖 pytesseract，请安装后重试") from exc
        try:
            image = Image.open(str(file_path))
            return pytesseract.image_to_string(image)
        except Exception as exc:
            raise ValueError(f'图片识别失败: {exc}') from exc

    raise ValueError('该资源类型不支持文本提取')

async def extract_resource_text_content(resource: TeachingResourceModel, max_chars: int) -> dict:
    # 如果资源有外部链接，使用联网功能提取
    if resource.external_url:
        provider = (getattr(settings, "AIGC_PROVIDER", None) or "deepseek").lower()
        if provider in {"kimi", "moonshot"}:
            try:
                client = AIGCClient()
                raw_text = await client.extract_url_content(resource.external_url, max_chars=max_chars)
                normalized = _normalize_text(raw_text)
                total_chars = len(normalized)
                truncated = total_chars > max_chars
                text = normalized[:max_chars] if truncated else normalized
                
                return {
                    "resource_id": resource.id,
                    "text": text,
                    "total_chars": total_chars,
                    "truncated": truncated,
                    "source": "external_url"
                }
            except Exception as exc:
                raise HTTPException(status_code=500, detail=f"外部链接内容提取失败: {exc}") from exc
        else:
            raise HTTPException(status_code=400, detail="外部链接内容提取仅支持 Kimi 提供商")
    
    # 如果没有文件路径，返回错误
    if not resource.file_path:
        raise HTTPException(status_code=400, detail="该资源没有可读取的文件或外部链接")

    file_path = Path(resource.file_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="资源文件不存在")

    # 检查文件大小，如果太大给出警告
    file_size = file_path.stat().st_size
    if file_size > 50 * 1024 * 1024:  # 50MB
        raise HTTPException(status_code=400, detail="文件过大（超过50MB），无法提取文本")

    provider = (getattr(settings, "AIGC_PROVIDER", None) or "deepseek").lower()
    if provider in {"kimi", "moonshot"}:
        try:
            client = AIGCClient()
            raw_text = await client.extract_file_text(str(file_path))
            raw_text = _normalize_kimi_content(raw_text)
            if _needs_image_fallback(raw_text) and file_path.suffix.lower() == ".docx":
                fallback_text = await _extract_docx_images_text(file_path, client)
                if fallback_text:
                    raw_text = fallback_text
        except httpx.TimeoutException:
            raise HTTPException(status_code=504, detail="文件处理超时，请稍后重试或尝试较小的文件")
        except httpx.HTTPStatusError as exc:
            # 处理 HTTP 状态错误
            status_code = exc.response.status_code
            if status_code in [502, 503, 504]:
                raise HTTPException(
                    status_code=503, 
                    detail=f"Kimi API 服务暂时不可用（{status_code}），请稍后重试。这是 Kimi 服务端的问题，已自动重试3次仍然失败。"
                )
            elif status_code == 429:
                raise HTTPException(status_code=429, detail="Kimi API 请求频率超限，请稍后重试")
            else:
                raise HTTPException(status_code=500, detail=f"Kimi API 返回错误 {status_code}: {exc}")
        except Exception as exc:
            error_msg = str(exc)
            if "timeout" in error_msg.lower():
                raise HTTPException(status_code=504, detail="文件处理超时，请稍后重试")
            raise HTTPException(status_code=500, detail=f"Kimi 文件解析失败: {exc}") from exc
    else:
        try:
            raw_text = await run_in_threadpool(
                _extract_text_sync,
                file_path,
                resource.file_format or '',
                resource.resource_type,
            )
        except ValueError as exc:
            raise HTTPException(status_code=400, detail=str(exc)) from exc
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"文本提取失败: {exc}") from exc

    normalized = _normalize_text(raw_text)
    total_chars = len(normalized)
    truncated = total_chars > max_chars
    text = normalized[:max_chars] if truncated else normalized

    return {
        "resource_id": resource.id,
        "text": text,
        "total_chars": total_chars,
        "truncated": truncated,
        "source": "file"
    }

resource_service = ResourceService()

@router.get("/", summary="获取教学资源列表")
async def get_resources(
    keyword: str = Query(None, description="关键词搜索"),
    resource_type: str = Query(None, description="资源类型"),
    software_engineering_chapter: str = Query(None, description="软件工程章节"),
    course_id: int = Query(None, description="课程ID"),
    chapter_id: int = Query(None, description="章节ID"),
    theme_category_id: int = Query(None, description="思政主题分类ID"),
    page: int = Query(1, ge=1, description="页码"),
    page_size: int = Query(10, ge=1, le=100, description="每页数量"),
    current_user: User = Depends(AuthControl.is_authed)
):
    search_request = ResourceSearchRequest(
        keyword=keyword,
        resource_type=resource_type,
        software_engineering_chapter=software_engineering_chapter,
        course_id=course_id,
        chapter_id=chapter_id,
        theme_category_id=theme_category_id,
        page=page,
        page_size=page_size
    )

    result = await resource_service.get_resources_with_search(search_request, current_user.id)

    # 转换为响应格式
    resources = [TeachingResource.from_orm(item) for item in result["items"]]

    return {
        "items": resources,
        "total": result["total"],
        "page": result["page"],
        "page_size": result["page_size"],
        "pages": result["pages"]
    }

@router.get("/{resource_id}/extract-text", summary="提取教学资源文本")
async def extract_resource_text(
    resource_id: int,
    max_chars: int = Query(8000, ge=100, le=100000, description="最大返回字符数（100-100000）"),
    current_user: User = Depends(AuthControl.is_authed)
):
    resource = await TeachingResourceModel.get_or_none(id=resource_id)
    if not resource:
        raise HTTPException(status_code=404, detail="资源不存在")

    # 检查权限
    if not resource.is_public and resource.uploader_id != current_user.id:
        raise HTTPException(status_code=403, detail="无权访问该资源")

    return await extract_resource_text_content(resource, max_chars)


@router.get("/recommended/list", summary="获取推荐教学资源")
async def get_recommended_resources(
    course_id: int = Query(None, description="课程ID"),
    chapter_id: int = Query(None, description="章节ID"),
    theme_category_id: int = Query(None, description="思政主题分类ID"),
    resource_types: str = Query(None, description="资源类型列表，逗号分隔"),
    limit: int = Query(5, ge=1, le=20, description="返回数量"),
    current_user: User = Depends(AuthControl.is_authed)
):
    types = [t.strip() for t in resource_types.split(",") if t.strip()] if resource_types else None
    resources = await RecommendationService.get_resource_recommendations(
        course_id=course_id,
        chapter_id=chapter_id,
        theme_category_id=theme_category_id,
        resource_types=types,
        limit=limit
    )
    return [TeachingResource.from_orm(item) for item in resources]

@router.post("/", summary="创建教学资源")
async def create_resource(
    title: str = Form(...),
    description: str = Form(None),
    resource_type: str = Form(...),
    external_url: str = Form(None),
    tags: str = Form(None),
    software_engineering_chapter: str = Form(None),
    course_id: int = Form(None),
    chapter_id: int = Form(None),
    theme_category_id: int = Form(None),
    is_public: bool = Form(True),
    file: UploadFile = File(None),
    current_user: User = Depends(AuthControl.is_authed)
):
    # 处理文件上传
    file_info = None
    if file and file.filename:
        file_info = await resource_service.upload_file(file, current_user.id)

    # 处理标签
    tags_list = []
    if tags:
        tags_list = [tag.strip() for tag in tags.split(",") if tag.strip()]

    resource_data = {
        "title": title,
        "description": description,
        "resource_type": resource_type,
        "external_url": external_url,
        "tags": tags_list,
        "software_engineering_chapter": software_engineering_chapter,
        "course_id": course_id,
        "chapter_id": chapter_id,
        "theme_category_id": theme_category_id,
        "is_public": is_public
    }

    # 如果有上传文件，添加文件相关信息
    if file_info:
        resource_data.update({
            "file_url": file_info["file_url"],
            "file_size": file_info["file_size"],
            "file_format": file_info["file_format"],
            "download_url": file_info["download_url"],
            "preview_url": file_info["preview_url"]
        })

    resource_in = TeachingResourceCreate(**resource_data)
    resource = await resource_service.create_resource(
        resource_in,
        current_user.id,
        file_info["file_path"] if file_info else None
    )

    return TeachingResource.from_orm(resource)


@router.post("/json", summary="创建教学资源（JSON，无文件上传）")
async def create_resource_json(
    resource_in: TeachingResourceCreate,
    current_user: User = Depends(AuthControl.is_authed)
):
    """
    允许前端以 JSON 方式创建教学资源（不含文件），避免表单缺少字段导致 422。
    如需上传文件，请使用表单上传接口。
    """
    resource = await resource_service.create_resource(resource_in, current_user.id)
    return TeachingResource.from_orm(resource)

@router.get("/{resource_id}", summary="获取教学资源详情")
async def get_resource(
    resource_id: int,
    current_user: User = Depends(AuthControl.is_authed)
):
    resource = await TeachingResourceModel.get_or_none(id=resource_id).prefetch_related('uploader')
    if not resource:
        raise HTTPException(status_code=404, detail="资源不存在")

    # 检查权限
    if not resource.is_public and resource.uploader_id != current_user.id:
        raise HTTPException(status_code=403, detail="无权访问该资源")

    # 增加使用次数
    await resource_service.update_resource_usage(resource_id)

    return TeachingResource.from_orm(resource)

@router.put("/{resource_id}", summary="更新教学资源")
async def update_resource(
    resource_id: int,
    resource_in: TeachingResourceUpdate,
    current_user: User = Depends(AuthControl.is_authed)
):
    resource = await TeachingResourceModel.get_or_none(id=resource_id)
    if not resource:
        raise HTTPException(status_code=404, detail="资源不存在")

    # 检查权限：只有上传者或超级管理员可以修改
    if resource.uploader_id != current_user.id and not current_user.is_superuser:
        raise HTTPException(status_code=403, detail="无权修改该资源")

    update_data = resource_in.dict(exclude_unset=True)
    update_data = await _hydrate_course_chapter(update_data)
    await resource.update_from_dict(update_data)
    await resource.save()

    return TeachingResource.from_orm(resource)

@router.delete("/{resource_id}", summary="删除教学资源")
async def delete_resource(
    resource_id: int,
    current_user: User = Depends(AuthControl.is_authed)
):
    resource = await TeachingResourceModel.get_or_none(id=resource_id)
    if not resource:
        raise HTTPException(status_code=404, detail="资源不存在")

    # 检查权限：只有上传者或超级管理员可以删除
    if resource.uploader_id != current_user.id and not current_user.is_superuser:
        raise HTTPException(status_code=403, detail="无权删除该资源")

    # 删除文件（如果存在）
    if resource.file_path and os.path.exists(resource.file_path):
        os.remove(resource.file_path)

    await resource.delete()
    return {"message": "资源删除成功"}

@router.post("/batch", summary="批量操作资源")
async def batch_operation(
    batch_request: BatchOperationRequest,
    current_user: User = Depends(AuthControl.is_authed)
):
    success_count = 0
    failed_count = 0
    failed_ids = []
    errors = []

    for resource_id in batch_request.target_ids:
        try:
            resource = await TeachingResourceModel.get_or_none(id=resource_id)
            if not resource:
                failed_count += 1
                failed_ids.append(resource_id)
                errors.append(f"资源 {resource_id} 不存在")
                continue

            # 检查权限
            if resource.uploader_id != current_user.id and not current_user.is_superuser:
                failed_count += 1
                failed_ids.append(resource_id)
                errors.append(f"无权操作资源 {resource_id}")
                continue

            # 执行批量操作
            if batch_request.operation == "delete":
                # 删除文件
                if resource.file_path and os.path.exists(resource.file_path):
                    os.remove(resource.file_path)
                await resource.delete()
            elif batch_request.operation == "public":
                await resource.update_from_dict({"is_public": True})
                await resource.save()
            elif batch_request.operation == "private":
                await resource.update_from_dict({"is_public": False})
                await resource.save()
            else:
                failed_count += 1
                failed_ids.append(resource_id)
                errors.append(f"不支持的操作类型: {batch_request.operation}")
                continue

            success_count += 1

        except Exception as e:
            failed_count += 1
            failed_ids.append(resource_id)
            errors.append(f"操作资源 {resource_id} 时出错: {str(e)}")

    return BatchOperationResponse(
        success_count=success_count,
        failed_count=failed_count,
        failed_ids=failed_ids,
        errors=errors
    )

@router.get("/hot/list", summary="获取热门资源")
async def get_hot_resources(
    limit: int = Query(10, ge=1, le=50, description="获取数量"),
    current_user: User = Depends(AuthControl.is_authed)
):
    resources = await resource_service.get_hot_resources(limit)
    return [TeachingResource.from_orm(resource) for resource in resources]

@router.get("/statistics/mine", summary="获取我的资源统计")
async def get_my_resources_statistics(
    current_user: User = Depends(AuthControl.is_authed)
):
    """获取当前用户的资源统计信息"""
    total = await TeachingResourceModel.filter(uploader_id=current_user.id).count()
    public = await TeachingResourceModel.filter(uploader_id=current_user.id, is_public=True).count()
    private = await TeachingResourceModel.filter(uploader_id=current_user.id, is_public=False).count()
    
    # 按类型统计
    resource_types = await TeachingResourceModel.filter(uploader_id=current_user.id).values_list("resource_type", flat=True)
    type_counts = {}
    for rt in resource_types:
        type_counts[rt] = type_counts.get(rt, 0) + 1
    
    return {
        "total": total,
        "public": public,
        "private": private,
        "by_type": type_counts
    }

@router.get("/download/{file_uuid}", summary="下载文件")
async def download_file(
    file_uuid: str,
    current_user: User = Depends(AuthControl.is_authed)
):
    # 根据文件UUID查找资源
    resource = await TeachingResourceModel.get_or_none(file_url__contains=file_uuid)
    if not resource:
        raise HTTPException(status_code=404, detail="文件不存在")

    # 检查权限
    if not resource.is_public and resource.uploader_id != current_user.id:
        raise HTTPException(status_code=403, detail="无权下载该文件")

    if not resource.file_path or not os.path.exists(resource.file_path):
        raise HTTPException(status_code=404, detail="文件不存在")

    # 流式返回文件
    def iterfile():
        with open(resource.file_path, mode="rb") as file_like:
            yield from file_like

    return StreamingResponse(
        iterfile(),
        media_type="application/octet-stream",
        headers={"Content-Disposition": f"attachment; filename={resource.title}"}
    )

@router.get("/types/list", summary="获取资源类型列表")
async def get_resource_types(
    current_user: User = Depends(AuthControl.is_authed)
):
    types = [
        {"value": "document", "label": "文档"},
        {"value": "video", "label": "视频"},
        {"value": "audio", "label": "音频"},
        {"value": "image", "label": "图片"},
        {"value": "presentation", "label": "演示文稿"},
        {"value": "simulation", "label": "虚拟仿真"},
        {"value": "link", "label": "外部链接"},
        {"value": "other", "label": "其他"}
    ]
    return types

@router.get("/themes/list", summary="获取思政主题列表")
async def get_ideological_themes(
    current_user: User = Depends(AuthControl.is_authed)
):
    # 返回常见思政主题
    themes = [
        "工匠精神",
        "创新精神",
        "团队协作",
        "责任担当",
        "诚信品质",
        "法治意识",
        "科学精神",
        "人文素养",
        "家国情怀",
        "国际视野"
    ]
    return themes
